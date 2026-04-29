from __future__ import annotations

from pathlib import Path

import re
import qrcode

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

THEME = {
    # Dark brand
    "bg_dark": RGBColor(10, 16, 30),
    "bg_dark_2": RGBColor(18, 30, 58),
    # Light canvas (blue-tinted so it never looks "white")
    "bg_light": RGBColor(234, 244, 255),
    "bg_light_2": RGBColor(214, 232, 255),
    # Surfaces
    "card": RGBColor(255, 255, 255),
    "card_line": RGBColor(226, 232, 241),
    "card_dark": RGBColor(18, 28, 50),
    "card_dark_line": RGBColor(60, 80, 120),
    # Text
    "text": RGBColor(20, 28, 40),
    "muted": RGBColor(84, 98, 122),
    "ink": RGBColor(242, 246, 255),  # on dark
    # Accents
    "accent": RGBColor(23, 130, 255),  # vivid blue
    "accent2": RGBColor(255, 176, 0),  # amber
    "ok": RGBColor(38, 191, 120),
    "danger": RGBColor(230, 67, 88),
}

# Modern default on Windows (Arabic supported)
FONT_AR = "Segoe UI"

RLM = "\u200f"  # Right-to-left mark
LRM = "\u200e"  # Left-to-right mark


def rtl(s: str) -> str:
    return f"{RLM}{s}{RLM}"


def ltr(s: str) -> str:
    return f"{LRM}{s}{LRM}"


def _norm_text(s: str) -> str:
    """
    Normalize spacing for projector-friendly Arabic slides.
    - Ensure clean spacing around the long dash separator (—).
    - Collapse duplicate whitespace.
    """
    if not s:
        return s
    s = s.replace("\u00a0", " ")
    # Keep numeric ranges like 10–18 without spaces, but enforce spacing for separators.
    # Also wrap the separator with RLM to avoid RTL rendering "sticking" to adjacent words.
    sep = f" {RLM}—{RLM} "
    s = re.sub(r"\s*[—–]\s*", sep, s)  # common separators
    s = re.sub(r"(?<!\d)\s*-\s*(?!\d)", sep, s)  # hyphen only when not a numeric range
    s = re.sub(r"[ \t]{2,}", " ", s).strip()
    return s


def _rtlize_text_frame(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
    except Exception:
        pass
    for p in tf.paragraphs:
        _set_rtl(p)


def _set_rtl(paragraph):
    # Enforce RTL in OOXML + right alignment.
    paragraph.alignment = PP_ALIGN.RIGHT
    try:
        ppr = paragraph._p.get_or_add_pPr()
        # a:pPr has an 'rtl' attribute in OOXML
        ppr.set("rtl", "1")
    except Exception:
        pass


def _set_font(run, *, name="Arial", size=20, bold=False, color=THEME["text"]):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Hint Arabic language for better shaping/RTL behavior
    try:
        rpr = run._r.get_or_add_rPr()
        rpr.set(qn("a:lang"), "ar-SA")
    except Exception:
        pass


def _add_link_run(paragraph, text: str, url: str, *, size=18, bold=False, color=None):
    r = paragraph.add_run()
    r.text = _norm_text(text)
    _set_font(r, name=FONT_AR, size=size, bold=bold, color=color or THEME["accent"])
    r.font.underline = True
    r.hyperlink.address = url
    return r


def _rounded_panel(slide, left, top, width, height, *, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    return shp


def _badge(slide, text: str, *, left, top, fill=THEME["accent2"], fg=THEME["bg_dark"]):
    b = _rounded_panel(slide, left, top, Inches(1.7), Inches(0.42), fill_rgb=fill, line_rgb=None)
    tx = slide.shapes.add_textbox(b.left, b.top + Inches(0.02), b.width, b.height)
    tf = tx.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(text)
    _set_font(r, name=FONT_AR, size=13, bold=True, color=fg)


def add_cover(slide, *, title: str, subtitle: str, hero_image: Path | None = None):
    # Full-bleed hero image + dark overlay + title block
    if hero_image and hero_image.exists():
        slide.shapes.add_picture(str(hero_image), 0, 0, width=SLIDE_W, height=SLIDE_H)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = THEME["bg_dark"]
        overlay.fill.fore_color.transparency = 0.18
        overlay.line.fill.background()
    else:
        add_bg(slide, variant="dark")

    # Title block
    block = _rounded_panel(
        slide,
        Inches(1.05),
        Inches(1.0),
        Inches(10.8),
        Inches(3.0),
        fill_rgb=RGBColor(18, 28, 50),
        line_rgb=RGBColor(60, 80, 120),
    )
    block.fill.fore_color.transparency = 0.08

    tx = slide.shapes.add_textbox(block.left + Inches(0.55), block.top + Inches(0.35), block.width - Inches(1.1), Inches(2.2))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    _rtlize_text_frame(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=46, bold=True, color=THEME["ink"])

    p2 = tf.add_paragraph()
    _set_rtl(p2)
    r2 = p2.add_run()
    r2.text = _norm_text(subtitle)
    _set_font(r2, name=FONT_AR, size=20, bold=False, color=RGBColor(190, 205, 230))

    _badge(slide, "Pitch Deck", left=Inches(1.05), top=Inches(4.15))



def add_bg(slide, *, variant: str = "light"):
    # Set a real slide background first (always visible).
    if variant == "dark":
        c1, c2 = THEME["bg_dark"], THEME["bg_dark_2"]
    else:
        c1, c2 = THEME["bg_light"], THEME["bg_light_2"]

    bgf = slide.background.fill
    bgf.solid()
    bgf.fore_color.rgb = c1

    # Add a subtle "gradient-like" second tone using EMU-safe integers.
    # Avoid float math on Length objects (can render as 0 in some viewers).
    w = Emu(int(SLIDE_W))
    h = Emu(int(SLIDE_H))
    split = Emu(int(int(SLIDE_H) * 58 // 100))
    bottom_h = Emu(int(SLIDE_H) - int(split))

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, split)
    top.fill.solid()
    top.fill.fore_color.rgb = c1
    top.line.fill.background()

    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, split, w, bottom_h)
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = c2
    bottom.line.fill.background()

    # Brand header for light slides so they never look "plain white".
    if variant != "dark":
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, Emu(int(Inches(0.12))))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = THEME["accent"]
        hdr.line.fill.background()


def add_section(slide, title: str, subtitle: str | None = None):
    add_bg(slide, variant="dark")

    tx = slide.shapes.add_textbox(Inches(1.05), Inches(2.3), Inches(11.5), Inches(2.4))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    _rtlize_text_frame(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=48, bold=True, color=THEME["ink"])

    if subtitle:
        p2 = tf.add_paragraph()
        _set_rtl(p2)
        r2 = p2.add_run()
        r2.text = _norm_text(subtitle)
        _set_font(r2, name=FONT_AR, size=20, bold=False, color=RGBColor(190, 205, 230))


def add_agenda(slide, items: list[str]):
    add_bg(slide, variant="dark")
    card = _rounded_panel(slide, Inches(0.9), Inches(1.1), Inches(11.6), Inches(5.8), fill_rgb=THEME["card_dark"], line_rgb=THEME["card_dark_line"])

    title = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(0.45), card.width - Inches(1.2), Inches(0.7))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = "المحاور"
    _set_font(r, name=FONT_AR, size=36, bold=True, color=THEME["ink"])

    body = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(1.35), card.width - Inches(1.2), card.height - Inches(2.0))
    tfb = body.text_frame
    tfb.clear()
    tfb.word_wrap = True
    for i, item in enumerate(items):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        _set_rtl(p)
        p.text = f"• {_norm_text(item)}"
        p.runs[0].font.name = FONT_AR
        p.runs[0].font.size = Pt(24)
        p.runs[0].font.color.rgb = RGBColor(200, 210, 225)


def _ensure_qr(path: Path, data: str):
    path.parent.mkdir(parents=True, exist_ok=True)
    img = qrcode.make(data)
    img.save(path)


def add_contact(slide, *, phone: str, whatsapp_url: str, email: str):
    add_bg(slide, variant="dark")

    card = _rounded_panel(slide, Inches(0.9), Inches(1.1), Inches(11.6), Inches(5.8), fill_rgb=THEME["card_dark"], line_rgb=THEME["card_dark_line"])

    # Title
    t = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(0.45), Inches(7.0), Inches(0.7))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = "التواصل"
    _set_font(r, name=FONT_AR, size=36, bold=True, color=THEME["ink"])

    # Links list
    links = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(1.35), Inches(7.0), Inches(3.6))
    tfl = links.text_frame
    tfl.clear()
    tfl.word_wrap = True

    p1 = tfl.paragraphs[0]
    _set_rtl(p1)
    r1 = p1.add_run()
    r1.text = "واتساب: "
    _set_font(r1, name=FONT_AR, size=24, bold=False, color=RGBColor(200, 210, 225))
    _add_link_run(p1, "فتح المحادثة", whatsapp_url, size=24, bold=True)

    p2 = tfl.add_paragraph()
    _set_rtl(p2)
    r2 = p2.add_run()
    r2.text = "اتصال: "
    _set_font(r2, name=FONT_AR, size=24, bold=False, color=RGBColor(200, 210, 225))
    _add_link_run(p2, phone, f"tel:{phone.replace(' ', '')}", size=24, bold=True)

    p3 = tfl.add_paragraph()
    _set_rtl(p3)
    r3 = p3.add_run()
    r3.text = "بريد: "
    _set_font(r3, name=FONT_AR, size=24, bold=False, color=RGBColor(200, 210, 225))
    _add_link_run(p3, email, f"mailto:{email}", size=24, bold=True)

    # QR
    qr_path = ROOT / "assets" / "qr" / "whatsapp.png"
    _ensure_qr(qr_path, whatsapp_url)
    qr_panel = _rounded_panel(slide, card.left + Inches(8.15), card.top + Inches(1.2), Inches(3.15), Inches(3.15), fill_rgb=RGBColor(20, 32, 62), line_rgb=THEME["card_dark_line"])
    slide.shapes.add_picture(str(qr_path), qr_panel.left + Inches(0.3), qr_panel.top + Inches(0.3), width=Inches(2.55), height=Inches(2.55))
    cap = slide.shapes.add_textbox(qr_panel.left, qr_panel.top + Inches(2.75), qr_panel.width, Inches(0.5))
    tfc = cap.text_frame
    tfc.clear()
    tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    pc = tfc.paragraphs[0]
    _set_rtl(pc)
    rc = pc.add_run()
    rc.text = "امسح الـ QR لواتساب"
    _set_font(rc, name=FONT_AR, size=15, bold=True, color=RGBColor(200, 210, 225))


def add_case_study(slide, title: str, bullets: list[str]):
    add_bg(slide, variant="dark")
    card = _rounded_panel(slide, Inches(0.9), Inches(1.1), Inches(11.6), Inches(5.8), fill_rgb=THEME["card_dark"], line_rgb=THEME["card_dark_line"])

    tx = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(0.45), card.width - Inches(1.2), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    body = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(1.3), card.width - Inches(1.2), card.height - Inches(1.8))
    tfb = body.text_frame
    tfb.clear()
    tfb.word_wrap = True
    for i, b in enumerate(bullets[:8]):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        _set_rtl(p)
        p.text = f"• {_norm_text(b)}"
        p.runs[0].font.name = FONT_AR
        p.runs[0].font.size = Pt(24)
        p.runs[0].font.color.rgb = RGBColor(200, 210, 225)


def add_operating_model(slide, title: str, steps: list[str]):
    add_bg(slide, variant="dark")
    card = _rounded_panel(slide, Inches(0.9), Inches(1.1), Inches(11.6), Inches(5.8), fill_rgb=THEME["card_dark"], line_rgb=THEME["card_dark_line"])

    tx = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(0.45), card.width - Inches(1.2), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    x0 = card.left + Inches(0.6)
    y0 = card.top + Inches(1.35)
    w = card.width - Inches(1.2)
    h = card.height - Inches(2.0)
    rows = 3
    cols = 2
    pad = Inches(0.25)
    cell_w = (w - pad) / cols
    cell_h = (h - pad * (rows - 1)) / rows

    for idx, step in enumerate(steps[:6]):
        x = x0 + (idx % cols) * (cell_w + pad)
        y = y0 + (idx // cols) * (cell_h + pad)
        box = _rounded_panel(slide, x, y, cell_w, cell_h, fill_rgb=RGBColor(20, 32, 62), line_rgb=THEME["card_dark_line"])
        t = slide.shapes.add_textbox(box.left + Inches(0.25), box.top + Inches(0.12), box.width - Inches(0.5), box.height - Inches(0.24))
        tfb = t.text_frame
        tfb.clear()
        tfb.word_wrap = True
        p = tfb.paragraphs[0]
        _set_rtl(p)
        r = p.add_run()
        r.text = _norm_text(step)
        _set_font(r, name=FONT_AR, size=20, bold=True, color=THEME["ink"])


def add_deal_structure(slide):
    add_bg(slide, variant="dark")

    card = _rounded_panel(
        slide,
        Inches(0.9),
        Inches(1.1),
        Inches(11.6),
        Inches(5.8),
        fill_rgb=THEME["card_dark"],
        line_rgb=THEME["card_dark_line"],
    )

    t = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(0.45), card.width - Inches(1.2), Inches(0.7))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = "هيكل الصفقة وتدفق الأموال"
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    x0 = card.left + Inches(0.6)
    y0 = card.top + Inches(1.35)
    w0 = Inches(5.3)
    h0 = Inches(1.15)
    gap = Inches(0.25)

    opts = [
        ("استحواذ كامل", "شراء الأصل المتعثر بالكامل ثم إكماله وطرحه"),
        ("شراء حصص", "الدخول بحصة مؤثرة مع إدارة الإكمال والتسويق"),
        ("شراكة إنقاذ", "اتفاقية إكمال مقابل حصة أو هامش مع ضمانات"),
    ]
    for i, (hd, bd) in enumerate(opts):
        box = _rounded_panel(slide, x0, y0 + i * (h0 + gap), w0, h0, fill_rgb=RGBColor(20, 32, 62), line_rgb=THEME["card_dark_line"])
        tx = slide.shapes.add_textbox(box.left + Inches(0.35), box.top + Inches(0.18), box.width - Inches(0.7), box.height - Inches(0.36))
        tfb = tx.text_frame
        tfb.clear()
        tfb.word_wrap = True
        p1 = tfb.paragraphs[0]
        _set_rtl(p1)
        r1 = p1.add_run()
        r1.text = _norm_text(hd)
        _set_font(r1, name=FONT_AR, size=20, bold=True, color=THEME["ink"])
        p2 = tfb.add_paragraph()
        _set_rtl(p2)
        r2 = p2.add_run()
        r2.text = _norm_text(bd)
        _set_font(r2, name=FONT_AR, size=15, bold=False, color=RGBColor(200, 210, 225))

    fx = card.left + Inches(6.25)
    fy = card.top + Inches(1.35)
    fw = card.width - Inches(6.85)
    step_h = Inches(0.75)
    fgap = Inches(0.22)
    flow = [
        "تمويل أو مساهمات",
        "حساب ضمان",
        "صرف على الموقع",
        "بيع أو تأجير",
        "توزيع عوائد",
    ]
    for i, txt in enumerate(flow):
        b = _rounded_panel(slide, fx, fy + i * (step_h + fgap), fw, step_h, fill_rgb=RGBColor(18, 28, 50), line_rgb=THEME["card_dark_line"])
        tt = slide.shapes.add_textbox(b.left + Inches(0.25), b.top + Inches(0.12), b.width - Inches(0.5), b.height - Inches(0.24))
        tfc = tt.text_frame
        tfc.clear()
        pc = tfc.paragraphs[0]
        _set_rtl(pc)
        rc = pc.add_run()
        rc.text = _norm_text(txt)
        _set_font(rc, name=FONT_AR, size=18, bold=True, color=THEME["ink"])

        if i < len(flow) - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, b.left + b.width / 2 - Inches(0.15), b.top + b.height, Inches(0.3), fgap)
            arr.fill.solid()
            arr.fill.fore_color.rgb = THEME["accent"]
            arr.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None):
    add_bg(slide, variant="dark")

    tx = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(2.0))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    _rtlize_text_frame(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=42, bold=True, color=THEME["ink"])

    if subtitle:
        p2 = tf.add_paragraph()
        _set_rtl(p2)
        r2 = p2.add_run()
        r2.text = _norm_text(subtitle)
        _set_font(r2, name=FONT_AR, size=20, bold=False, color=RGBColor(185, 200, 220))


def add_card_bullets(slide, header: str, bullets: list[str], *, left=0.8, top=1.4, w=12.0, h=5.3):
    add_bg(slide, variant="dark")
    # Soft shadow panel
    shadow = _rounded_panel(
        slide,
        Inches(left + 0.05),
        Inches(top + 0.08),
        Inches(w),
        Inches(h),
        fill_rgb=RGBColor(0, 0, 0),
        line_rgb=None,
    )
    shadow.fill.fore_color.transparency = 0.88

    card = _rounded_panel(
        slide,
        Inches(left),
        Inches(top),
        Inches(w),
        Inches(h),
        fill_rgb=THEME["card_dark"],
        line_rgb=THEME["card_dark_line"],
    )

    # Header
    hdr = slide.shapes.add_textbox(Inches(left + 0.6), Inches(top + 0.35), Inches(w - 1.2), Inches(0.8))
    tfh = hdr.text_frame
    tfh.clear()
    _rtlize_text_frame(tfh)
    p = tfh.paragraphs[0]
    r = p.add_run()
    r.text = _norm_text(header)
    _set_font(r, name=FONT_AR, size=30, bold=True, color=THEME["ink"])

    # Bullets
    body = slide.shapes.add_textbox(Inches(left + 0.6), Inches(top + 1.25), Inches(w - 1.2), Inches(h - 1.7))
    tf = body.text_frame
    tf.clear()
    _rtlize_text_frame(tf)
    for i, b in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_rtl(p)
        p.level = 0
        p.text = f"• {_norm_text(b)}"
        p.runs[0].font.name = FONT_AR
        p.runs[0].font.size = Pt(26)
        p.runs[0].font.color.rgb = RGBColor(200, 210, 225)


def add_image_message_slide(slide, title: str, bullets: list[str], image_paths: list[Path]):
    add_bg(slide, variant="dark")

    def _is_video_frame(p: Path) -> bool:
        return "video-frames" in str(p).replace("\\", "/")

    # Left text card
    card = _rounded_panel(slide, Inches(0.9), Inches(1.2), Inches(6.6), Inches(5.9), fill_rgb=THEME["card_dark"], line_rgb=THEME["card_dark_line"])

    tx = slide.shapes.add_textbox(card.left + Inches(0.55), card.top + Inches(0.45), card.width - Inches(1.1), Inches(1.0))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=28, bold=True, color=THEME["ink"])

    body = slide.shapes.add_textbox(card.left + Inches(0.55), card.top + Inches(1.35), card.width - Inches(1.1), card.height - Inches(1.9))
    tfb = body.text_frame
    tfb.clear()
    tfb.word_wrap = True
    for i, b in enumerate(bullets[:5]):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        _set_rtl(p)
        p.text = f"• {_norm_text(b)}"
        p.runs[0].font.name = FONT_AR
        p.runs[0].font.size = Pt(22)
        p.runs[0].font.color.rgb = RGBColor(200, 210, 225)

    # Right images panel
    panel = _rounded_panel(slide, Inches(7.75), Inches(1.2), Inches(4.6), Inches(5.9), fill_rgb=RGBColor(20, 32, 62), line_rgb=THEME["card_dark_line"])
    cols = 1
    rows = 2
    pad = Inches(0.25)
    cell_w = panel.width - pad * 2
    cell_h = (panel.height - pad * (rows + 1)) / rows

    for idx, img in enumerate(image_paths[:2]):
        if not img.exists():
            continue
        x = panel.left + pad
        y = panel.top + pad + idx * (cell_h + pad)
        pic = slide.shapes.add_picture(str(img), x, y, width=cell_w, height=cell_h)
        if _is_video_frame(img):
            # Crop to emphasize "construction" area (top of the WhatsApp frame)
            pic.crop_bottom = 0.45


def add_construction_gallery(slide, title: str, image_paths: list[Path]):
    add_bg(slide, variant="dark")
    card = _rounded_panel(slide, Inches(0.9), Inches(1.1), Inches(11.6), Inches(5.8), fill_rgb=THEME["card_dark"], line_rgb=THEME["card_dark_line"])

    tx = slide.shapes.add_textbox(card.left + Inches(0.6), card.top + Inches(0.45), card.width - Inches(1.2), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(title)
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    # 2x2 images
    x0 = card.left + Inches(0.6)
    y0 = card.top + Inches(1.25)
    w = card.width - Inches(1.2)
    h = card.height - Inches(1.75)
    pad = Inches(0.25)
    cell_w = (w - pad) / 2
    cell_h = (h - pad) / 2

    for idx, img in enumerate(image_paths[:4]):
        if not img.exists():
            continue
        x = x0 + (idx % 2) * (cell_w + pad)
        y = y0 + (idx // 2) * (cell_h + pad)
        slide.shapes.add_picture(str(img), x, y, width=cell_w, height=cell_h)


def add_kpis(slide, header: str, kpis: list[tuple[str, str]]):
    add_bg(slide, variant="dark")
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(12.0), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(header)
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    left = 0.8
    top = 1.7
    box_w = 3.9
    box_h = 1.5
    gap = 0.35

    for idx, (value, label) in enumerate(kpis):
        x = left + (idx % 3) * (box_w + gap)
        y = top + (idx // 3) * (box_h + gap)
        box = _rounded_panel(
            slide,
            Inches(x),
            Inches(y),
            Inches(box_w),
            Inches(box_h),
            fill_rgb=RGBColor(18, 28, 50),
            line_rgb=RGBColor(60, 80, 120),
        )

        tx = slide.shapes.add_textbox(Inches(x + 0.35), Inches(y + 0.20), Inches(box_w - 0.7), Inches(box_h - 0.35))
        tfb = tx.text_frame
        tfb.clear()
        p1 = tfb.paragraphs[0]
        _set_rtl(p1)
        r1 = p1.add_run()
        r1.text = _norm_text(value)
        _set_font(r1, name=FONT_AR, size=32, bold=True, color=THEME["accent"])

        p2 = tfb.add_paragraph()
        _set_rtl(p2)
        r2 = p2.add_run()
        r2.text = _norm_text(label)
        _set_font(r2, name=FONT_AR, size=15, bold=False, color=RGBColor(200, 210, 225))


def add_timeline(slide, header: str, steps: list[tuple[str, str]]):
    add_bg(slide, variant="dark")
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(12.0), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_rtl(p)
    r = p.add_run()
    r.text = _norm_text(header)
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    x0 = Inches(1.0)
    y0 = Inches(2.0)
    w = Inches(11.6)
    h = Inches(0.12)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y0 + Inches(1.4), w, h)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(60, 80, 120)
    line.line.fill.background()

    n = len(steps)
    for i, (when, what) in enumerate(steps):
        cx = x0 + (w * i / max(1, (n - 1)))
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), y0 + Inches(1.33), Inches(0.24), Inches(0.24))
        dot.fill.solid()
        dot.fill.fore_color.rgb = THEME["accent2"] if i == 0 else THEME["accent"]
        dot.line.fill.background()

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cx - Inches(1.8),
            y0 + (Inches(0.0) if i % 2 == 0 else Inches(1.9)),
            Inches(3.6),
            Inches(1.35),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(20, 30, 52)
        card.line.color.rgb = RGBColor(60, 80, 120)

        tx = slide.shapes.add_textbox(card.left + Inches(0.25), card.top + Inches(0.12), card.width - Inches(0.5), card.height - Inches(0.24))
        tfc = tx.text_frame
        tfc.clear()
        tfc.word_wrap = True
        p1 = tfc.paragraphs[0]
        _set_rtl(p1)
        r1 = p1.add_run()
        r1.text = _norm_text(when)
        _set_font(r1, name=FONT_AR, size=14, bold=True, color=THEME["accent2"])

        p2 = tfc.add_paragraph()
        _set_rtl(p2)
        r2 = p2.add_run()
        r2.text = _norm_text(what)
        _set_font(r2, name=FONT_AR, size=15, bold=False, color=RGBColor(215, 225, 238))


def add_video_support_slide(slide, header: str, notes: list[str], image_paths: list[Path]):
    add_bg(slide, variant="dark")
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(12.0), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    _rtlize_text_frame(tf)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _norm_text(header)
    _set_font(r, name=FONT_AR, size=32, bold=True, color=THEME["ink"])

    # Left: bullets
    shadow = _rounded_panel(slide, Inches(0.85), Inches(1.78), Inches(6.1), Inches(5.1), fill_rgb=RGBColor(0, 0, 0), line_rgb=None)
    shadow.fill.fore_color.transparency = 0.88
    card = _rounded_panel(slide, Inches(0.8), Inches(1.7), Inches(6.1), Inches(5.1), fill_rgb=RGBColor(18, 28, 50), line_rgb=RGBColor(60, 80, 120))

    body = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(5.3), Inches(4.3))
    tfb = body.text_frame
    tfb.clear()
    _rtlize_text_frame(tfb)
    for i, b in enumerate(notes):
        p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        _set_rtl(p)
        p.text = f"• {_norm_text(b)}"
        if p.runs:
            _set_font(p.runs[0], name=FONT_AR, size=22, bold=False, color=RGBColor(200, 210, 225))

    # Add a real clickable CTA link (WhatsApp) under the bullets
    cta = tfb.add_paragraph()
    _set_rtl(cta)
    cta.space_before = Pt(8)
    _add_link_run(
        cta,
        "تواصل معنا الآن عبر واتساب",
        "https://wa.me/966581125715",
        size=20,
        bold=True,
        color=THEME["accent"],
    )

    # Right: images grid
    x = Inches(7.2)
    y = Inches(1.7)
    w = Inches(5.2)
    h = Inches(5.1)
    panel = _rounded_panel(slide, x, y, w, h, fill_rgb=RGBColor(18, 28, 50), line_rgb=RGBColor(60, 80, 120))

    cols = 2
    rows = 2
    pad = Inches(0.25)
    cell_w = (w - pad * (cols + 1)) / cols
    cell_h = (h - pad * (rows + 1)) / rows

    for idx, img in enumerate(image_paths[:4]):
        if not img.exists():
            continue
        cx = x + pad + (idx % cols) * (cell_w + pad)
        cy = y + pad + (idx // cols) * (cell_h + pad)
        slide.shapes.add_picture(str(img), cx, cy, width=cell_w, height=cell_h)


def build(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W  # 16:9 widescreen
    prs.slide_height = SLIDE_H

    # Slides
    frames_dir = ROOT / "assets" / "video-frames"
    hero = frames_dir / "frame_0018.jpg"
    add_cover(
        prs.slides.add_slide(prs.slide_layouts[6]),
        title="الاستثمار الآمن في التطوير العقاري النوعي",
        subtitle="نموذج مبتكر يركز على الاستحواذ الإنقاذي وإكمال المشاريع المتعثرة لتقليل المخاطر وتسريع دوران رأس المال — جدة",
        hero_image=hero if hero.exists() else None,
    )

    add_agenda(
        prs.slides.add_slide(prs.slide_layouts[6]),
        [
            "فرصة السوق",
            "نبذة عن الشركة",
            "مفهوم الشركة",
            "النموذج المالي والتشغيلي",
            "هيكل الصفقة",
            "الأرقام للسكني والتجاري وخارطة الطريق",
            "أسئلة المستثمر",
            "التواصل",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "نبذة عن الشركة — مساهمة مقفلة",
        [
            "شركة مساهمة مقفلة تُسجَّل رسميًا لدى وزارة التجارة وفق الأنظمة المعمول بها في المملكة",
            "هيكل أسهم واضح ورأس مال محدد — يسهّل الدخول والخروج ويعزز الانضباط المالي",
            "حوكمة مؤسسية عبر مجلس إدارة ومحاضر واجتماعات وتقارير دورية للمساهمين",
            "تمنح ثقة أعلى للشركاء والبنوك مقارنة بهياكل فردية أو غير منظمة",
            "قابلية التوسع لاحقًا — فتح الباب لإدراج مستقبلي عند نضج المحفظة",
        ],
        top=1.35,
        h=5.35,
    )

    add_section(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "فرصة السوق",
        "الفجوة + لماذا الآن؟",
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "المشكلة في السوق",
        [
            "مليارات الريالات محاصرة في هياكل خرسانية ومشاريع غير مكتملة لا تحقق عائدًا",
            "فجوة السيولة — ارتفاع أسعار الفائدة وتذبذب تكاليف مواد البناء",
            "تعثر مشاريع مكتملة بنسبة 40%–60% يخلق أصولًا مجمّدة تبحث عن مخرج مالي وفني",
            "تعثر الملاك بسبب ضائقة مالية/قانونية أو غياب التمويل اللازم للإكمال",
            "ضغط تنظيمي لتسريع الإنجاز يدفع المطور المتعثر للبيع أو الشراكة لتجنب العقوبات",
            "خسائر مزدوجة: خسارة المالك الأصلي + حرمان السوق من وحدات مطلوبة",
            "فجوة العرض والطلب في جدة ترفع الحاجة لحلول إنقاذ وإحياء الأصول",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "مفهوم الشركة",
        [
            "شركة مساهمة مغلقة متخصصة في «الاستحواذ الإنقاذي» والإكمال الإنشائي",
            "شراء حصص في المشاريع المتعثرة أو الاستحواذ الكامل عليها ثم إعادة هيكلة الميزانية",
            "ضخ السيولة لإتمام البناء بجداول زمنية مضغوطة وتحويل الأصل المتعثر إلى منتج جاهز",
            "الدخول بعد اكتمال الأساسات والهيكل الإنشائي لتقليل المخاطر وتسريع التنفيذ",
        ],
    )

    add_section(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "النموذج المالي والتشغيلي",
        "تشغيل + تمويل + إكمال",
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "لماذا الآن؟",
        [
            "أصول متعثرة/متوقفة تخلق فرص استحواذ بأسعار أفضل من السوق",
            "فجوة سيولة لدى بعض الملاك والمطورين تفتح باب الشراكات الإنقاذية",
            "سوق سكني/تجاري يحتاج وحدات جاهزة بسرعة بعد الإكمال",
            "تركيزنا: إكمال البناء لتسريع دوران رأس المال — وليس البنية التحتية",
        ],
    )

    add_operating_model(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "نموذج التشغيل",
        [
            "1) رصد فرصة",
            "2) عناية واجبة هندسية/قانونية",
            "3) هيكلة صفقة — شراء أو شراكة إنقاذ",
            "4) إكمال + ضبط جودة",
            "5) طرح للبيع/التأجير + تسويق",
            "6) خروج + توزيع أرباح",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "من نحن",
        [
            "استحواذ: رصد الفرص والتفاوض وهيكلة الصفقات",
            "هندسة: فحص الهيكل وتقدير نطاق الإكمال وضبط الجودة",
            "تشغيل وتسويق: إدارة المقاولين والتصريف عبر البيع أو التأجير",
            "شفافية: تقارير دورية وربط الصرف بالإنجاز",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "آلية اتخاذ القرار",
        [
            "تحديد فرصة مبدئية وفق معايير الاختيار",
            "عناية واجبة هندسية وقانونية ومالية",
            "عرض صفقة مختصر للجنة الاستثمار",
            "قرار: قبول أو رفض أو إعادة تسعير",
            "تفعيل حساب الضمان ثم بدء الإكمال",
        ],
    )

    add_deal_structure(prs.slides.add_slide(prs.slide_layouts[6]))

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "معايير اختيار الفرص",
        [
            "نسبة إنجاز مناسبة — هيكل قائم وقابلية إكمال خلال مدة قصيرة",
            "سلامة إنشائية مثبتة بفحص طرف ثالث",
            "وضع تراخيص واضح وقابل للإغلاق السريع",
            "موقع قابل للتصريف وطلب حقيقي — سكني أو تجاري",
            "ميزانية إكمال محددة + احتياطي مخاطر",
            "حد أدنى للعائد المستهدف حسب النوع — سكني أو تجاري",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "إدارة المخاطر",
        [
            "مخاطر هندسية: فحص + تقرير + نطاق أعمال واضح",
            "مخاطر قانونية: تدقيق ملكية/التزامات/مقاولين",
            "مخاطر التنفيذ: عقود جزائية + مراقبة جودة + جدول واقعي",
            "مخاطر التكلفة: BOQ + احتياطي + ضبط تغييرات",
            "مخاطر السوق: تسعير مرحلي + قنوات بيع/تأجير",
            "حوكمة: حساب ضمان وربط الصرف بنِسَب الإنجاز",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "الشراكات للمجمعات الجاهزة والبناء الذكي",
        [
            "التعاون مع الإسكان/شركات كبيرة لتسريع التسليم",
            "حلول البناء الذكي كقيمة مضافة — أمن وتحكم وطاقة وصيانة",
            "مواصفات معيارية قابلة للتكرار على عدة مشاريع",
            "نموذج تشغيل/صيانة مع شركاء مختصين بعد التسليم",
        ],
    )

    add_case_study(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "صفقة نموذجية — مثال تقديري",
        [
            "مشروع سكني متوقف بنسبة إنجاز 50% — هيكل قائم",
            "نستحوذ بخصم تقديري حسب المخاطر وحالة السوق",
            "إكمال خلال 6–12 شهر — تقديري حسب حجم المشروع",
            "طرح للبيع على مراحل لتقليل المخاطر وتسريع التدفق النقدي",
            "في التجاري: التركيز على مستأجر/عقد رئيسي لتقليل مخاطر الإشغال",
            "الأرقام الدقيقة تُحدد بعد العناية الواجبة لكل فرصة",
        ],
    )

    add_construction_gallery(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "نماذج مشاريع غير مكتملة — أمثلة مصوّرة",
        [
            ROOT / "assets" / "web-construction" / "c6.jpg",
            ROOT / "assets" / "web-construction" / "c9.jpg",
            ROOT / "assets" / "web-construction" / "c2.jpg",
            ROOT / "assets" / "web-construction" / "c4.jpg",
        ],
    )

    add_image_message_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "ما نستهدفه — المنتجات",
        [
            "عماير شقق سكنية",
            "فلل وملاحق",
            "مكاتب — تجاري",
            "تجزئة — محلات ووحدات تجارية",
            "لا نستهدف مشاريع البنية التحتية",
        ],
        [
            ROOT / "assets" / "web-construction" / "c6.jpg",
            ROOT / "assets" / "web-construction" / "c9.jpg",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "ماذا نفعل؟ — نطاق العمل",
        [
            "الاستحواذ على مشروع متعثر/غير مكتمل أو الدخول كشريك إنقاذ",
            "إكمال البناء — تشطيب وتمديدات وأنظمة للوصول إلى جاهزية التشغيل أو البيع",
            "تسريع دوران رأس المال عبر تقليص زمن الوصول للسوق",
            "حوكمة صرف التنفيذ — حساب ضمان وربط الدفع بالإنجاز",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "المحركات المالية والاستثمارية",
        [
            "الاستفادة من خصم الأصول المتعثرة عند الاستحواذ",
            "ملاءة شركة مساهمة تسهّل الوصول لتمويل بتكلفة أقل مقارنة بالمطورين الأفراد",
            "رفع ROI عبر شراء بسعر التكلفة/أقل ثم بيع بسعر السوق عند الجاهزية",
            "دوران أسرع لرأس المال نتيجة تقليص مراحل ما قبل التنفيذ",
        ],
    )

    add_image_message_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "البناء الذكي للمجمعات الجاهزة",
        [
            "إضافة حلول بناء ذكي للمجمعات السكنية الجاهزة عند الإكمال",
            "شراكات تنفيذ/تشغيل مع الإسكان أو الشركات الكبيرة",
            "رفع القيمة السوقية وتحسين التجربة — تحكم وأمن وطاقة",
            "جاهزية تسليم أسرع بمواصفات معيارية قابلة للتوسع",
        ],
        [
            ROOT / "assets" / "web-construction" / "c2.jpg",
            ROOT / "assets" / "web-construction" / "c4.jpg",
        ],
    )
    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "الاستراتيجية التشغيلية",
        [
            "تدقيق فني وهندسي كامل للمنشآت القائمة لضمان السلامة قبل ضخ السيولة",
            "تنفيذ متسارع باستخدام تقنيات البناء الحديثة",
            "حوكمة تشغيلية عبر حسابات ضمان لتوجيه الإنفاق إلى الموقع مباشرة",
            "إدارة مشاريع صارمة: تشطيب/تمديدات + تسويق + بيع للوصول لمنتج نهائي جاهز",
        ],
    )

    add_section(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "الأرقام",
        "KPIs تقديرية للسوق السعودي",
    )

    add_kpis(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "مؤشرات تقديرية — السكني",
        [
            (rtl("≈ 10–18%"), "عائد مستهدف — تقديري"),
            (f"≈ {ltr('18–36')} شهر", "مدة استرداد متوقعة"),
            (rtl("≈ 8–20%"), "خصم شراء أصل متعثر — تقديري"),
            (rtl("≈ 15–35%"), "تكلفة الإكمال من إجمالي قيمة المشروع"),
            (f"≈ {ltr('6–12')} شهر", "زمن التصريف بعد الإكمال — تقديري"),
            ("منخفضة–متوسطة", "مخاطر أقل بوجود أصل قائم — ليست صفرًا"),
        ],
    )

    add_kpis(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "مؤشرات تقديرية — التجاري",
        [
            (rtl("≈ 12–22%"), "عائد مستهدف — تقديري"),
            (f"≈ {ltr('24–48')} شهر", "مدة استرداد متوقعة"),
            (rtl("≈ 10–25%"), "خصم شراء أصل متعثر — تقديري"),
            (rtl("≈ 20–45%"), "تكلفة الإكمال من إجمالي قيمة المشروع"),
            (f"≈ {ltr('9–18')} شهر", "زمن التصريف/التأجير بعد الإكمال — تقديري"),
            ("متوسطة", "مخاطر أعلى نسبيًا — ارتباط بالإشغال والمستأجرين"),
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "أربع ركائز تجعل نموذجنا أقل مخاطرة",
        [
            "أصول ثابتة قائمة — أرض ومبنى قائم كضمان حقيقي لرأس المال",
            "اختصار الوقت بتجاوز مراحل الترخيص والحفر والأساسات",
            "دوران أسرع للمال: بيع خلال أشهر بدل 3–5 سنوات",
            "هوامش ربح أعلى بفضل الاستحواذ بسعر تنافسي",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "مقارنة مختصرة مع التطوير التقليدي",
        [
            "التقليدي: دخول من الصفر، مخاطر أعلى، ضمانات غالبًا أرض فقط",
            "التقليدي: استرداد 3–5 سنوات وهوامش عادةً أقل استقرارًا وترخيص يستغرق أشهر",
            "نموذجنا: مشروع قائم جزئيًا، ضمانات أرض + مبنى قائم، مخاطر منخفضة",
            "نموذجنا: السكني غالبًا أسرع في التصريف من التجاري، لكن التجاري قد يعتمد على عقود وتأجير",
            "نموذجنا: الاسترداد والهوامش تقديرية وتتغير حسب المدينة، المنتج، نسبة الإنجاز، والتسعير",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "لماذا جدة؟",
        [
            "سوق عقاري في مرحلة نمو، مع طلب سكني متصاعد حتى 2030",
            "مشاريع الرؤية واستثمارات البنية التحتية ترفع القيمة السوقية للأصول",
            "الطلب على الوحدات يتجاوز العرض المتاح → سرعة تصريف أعلى",
            "موقع استراتيجي على ساحل البحر الأحمر يدعم الاستثمار العقاري والسياحي",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "هيكل الشركة والحوكمة",
        [
            "شركة مساهمة مقفلة — مقرها جدة، المملكة العربية السعودية",
            "حوكمة مؤسسية وشفافية لحفظ حقوق المساهمين وتحقيق الأهداف الاستثمارية",
            "نطاق العمل: جدة — بناءً على الفرص المتاحة",
            "قيمة السهم: 23,500 ريال — حسب الوثيقة",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "مبررات التأسيس كمساهمة مغلقة + خطة الخروج",
        [
            "استدامة رأس المال: جمع سيولة ضخمة من مساهمين مؤسسين لضمان عدم التوقف",
            "مرونة وسرعة قرار الاستحواذ مقارنة بالشركات المساهمة العامة",
            "إمكانية الإدراج في سوق «نمو» مستقبلًا بعد بناء محفظة مشاريع ناجحة",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "الربط مع رؤية 2030",
        [
            "دعم مستهدفات الإسكان عبر إعادة ضخ الوحدات المعطلة إلى السوق",
            "معالجة المشاريع المتعثرة يحدّ من مخاطر هبوط السوق الناتج عن توقف المشاريع الكبيرة",
            "رفع كفاءة القطاع عبر تحويل الأصول المجمدة إلى منتجات سكنية جاهزة",
        ],
    )

    add_timeline(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "خارطة طريق حتى أول عائد — 18 شهرًا",
        [
            ("1–3 أشهر", "إتمام التأسيس القانوني، جمع رأس المال، تشكيل فريق التنفيذ"),
            ("3–6 أشهر", "مسح السوق، تحديد المشاريع المستهدفة، العناية الواجبة"),
            ("6–12 شهر", "الاستحواذ على أول مشروع وبدء أعمال الإكمال والتطوير"),
            ("12–18 شهر", "إطلاق المبيعات، تحقيق العائد، وتوزيع الأرباح بشفافية"),
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "أسئلة المستثمر المتوقعة",
        [
            "كيف نحمي رأس المال وما الضمانات؟",
            "كيف نحدد سعر الاستحواذ وميزانية الإكمال؟",
            "ما خطة الخروج في السكني مقابل التجاري؟",
            "كيف نخفف مخاطر التأخير وارتفاع التكلفة؟",
            "كيف نضمن الشفافية في الصرف والتقارير؟",
            "ما الحد الأقصى للمخاطر المقبولة لكل صفقة؟",
        ],
    )

    add_card_bullets(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "ما الذي نحتاجه الآن",
        [
            "فرص متعثرة جاهزة للتقييم داخل جدة",
            "شركاء تنفيذ وتشطيب بأنظمة جودة واضحة",
            "شركاء تسويق وبيع أو إدارة تأجير",
            "ممولون أو مساهمون لتسريع الإغلاق والإكمال",
        ],
    )

    # Video support slide (no identity claims — only what appears in overlay text)
    video_imgs = [
        frames_dir / "frame_0001.jpg",
        frames_dir / "frame_0018.jpg",
        frames_dir / "frame_0024.jpg",
        frames_dir / "frame_0030.jpg",
    ]
    add_section(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "مواد داعمة",
        "رسائل ولقطات من الفيديو",
    )
    add_video_support_slide(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "محتوى داعم من الفيديو — رسائل وعناوين ظاهرة",
        [
            "التركيز على تأسيس أو هيكلة شركة مساهمة مقفلة — وفق النص الظاهر في الفيديو",
            "ذكر «عدد الأسهم 2000 سهم» ضمن الرسالة التعريفية",
            "التأكيد على «فرص واعدة واستثمار بلا مخاطر» كرسالة تسويقية",
            "دعوة للتواصل والمشاركة — «شاركنا… تواصل معنا الآن»",
        ],
        video_imgs,
    )

    add_section(prs.slides.add_slide(prs.slide_layouts[6]), "التواصل", "روابط قابلة للنقر + QR")
    add_contact(
        prs.slides.add_slide(prs.slide_layouts[6]),
        phone="+966 58 112 5715",
        whatsapp_url="https://wa.me/966581125715",
        email="talalkenani@gmail.com",
    )

    add_title(
        prs.slides.add_slide(prs.slide_layouts[6]),
        "جاهزون للبدء",
        "انضم اليوم وكن جزءًا من قصة نجاح عقارية نوعية.",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main():
    out = ROOT / "exports" / "عرض_الاستثمار_الآمن_التطوير_العقاري_النوعي.pptx"
    build(out)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()

